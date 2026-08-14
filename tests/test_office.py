"""Tests for the office (.xlsx / .pptx) tools, incl. undoability of binary edits."""

import pytest

pytest.importorskip("openpyxl")
pytest.importorskip("pptx")
# NOTE: python-docx is gated per-test (see the read_docx tests), not module-wide,
# so the xlsx/pptx suite still runs when only python-docx is missing.

from opendot.reversibility.engine import Reversibility
from opendot.reversibility.snapshots import IgnoreRules
from opendot.tools.local import Toolbox


@pytest.fixture(autouse=True)
def isolated_store(tmp_path, monkeypatch):
    monkeypatch.setenv("OPENDOT_HOME", str(tmp_path / "store"))
    yield


def _tb(tmp_path):
    wd = tmp_path / "ws"
    wd.mkdir()
    rev = Reversibility(workdir=str(wd), rules=IgnoreRules())
    return Toolbox(str(wd), reversibility=rev, confirm=lambda p: True), wd, rev


def _make_xlsx(path):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "name"
    ws["B1"] = "score"
    ws["A2"] = "alice"
    ws["B2"] = 100
    wb.save(path)


def _add_xlsx_sheet(path):
    import openpyxl

    wb = openpyxl.load_workbook(path)
    ws = wb.create_sheet("Summary")
    ws["A1"] = "total"
    ws["B1"] = 100
    wb.save(path)


def _make_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Old Title"
    prs.save(path)


def _make_docx(path):
    from docx import Document

    doc = Document()
    doc.add_heading("Project Report", level=1)
    doc.add_paragraph("First paragraph of the body.")
    doc.add_paragraph("Second paragraph with details.")
    doc.save(path)


def test_office_tools_registered(tmp_path):
    tb, _, _ = _tb(tmp_path)
    names = {s["function"]["name"] for s in tb.specs()}
    assert {"read_xlsx", "edit_cell", "read_pptx", "edit_pptx_text"} <= names
    # read_docx is gated on python-docx being importable.
    import importlib.util

    if importlib.util.find_spec("docx") is not None:
        assert "read_docx" in names


def test_read_docx(tmp_path):
    pytest.importorskip("docx")
    tb, wd, _ = _tb(tmp_path)
    _make_docx(wd / "report.docx")

    out = tb.call("read_docx", {"path": "report.docx"})
    assert "Project Report" in out
    assert "First paragraph of the body." in out
    assert "Second paragraph with details." in out


def test_read_docx_missing_file(tmp_path):
    pytest.importorskip("docx")
    tb, _, _ = _tb(tmp_path)
    assert tb.call("read_docx", {"path": "nope.docx"}).startswith("error: file not found")


def test_read_and_edit_xlsx(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")

    out = tb.call("read_xlsx", {"path": "data.xlsx"})
    assert "alice" in out and "score" in out

    res = tb.call("edit_cell", {"path": "data.xlsx", "cell": "B2", "value": "120"})
    assert "100" in res and "120" in res

    import openpyxl

    wb = openpyxl.load_workbook(wd / "data.xlsx")
    assert wb.active["B2"].value == 120  # coerced to int


@pytest.mark.parametrize(
    ("value", "expected"),
    [
        ("7", 7),
        ("-5", -5),
        ("+7", 7),
        ("0", 0),
        ("3.14", 3.14),
        ("1e3", 1000.0),
        (".5", 0.5),
        # These must stay TEXT (the bug this fixes):
        ("007", "007"),  # leading-zero ID/ZIP/SKU
        ("1_000", "1_000"),  # underscore separators
        ("inf", "inf"),  # non-finite
        ("nan", "nan"),
        ("-inf", "-inf"),
        ("١٢٣", "١٢٣"),  # unicode digits
        ("2,000", "2,000"),  # thousands separator
        ("abc", "abc"),
        ("007-88-1234", "007-88-1234"),
        ("=A1+A2", "=A1+A2"),  # formula passes through
    ],
)
def test_coerce_cell_value(value, expected):
    from opendot.tools.office import _coerce_cell_value

    result = _coerce_cell_value(value)
    assert result == expected
    assert type(result) is type(expected)  # int vs float vs str must match


def test_edit_cell_preserves_leading_zeros_by_default(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    tb.call("edit_cell", {"path": "data.xlsx", "cell": "A2", "value": "007"})

    import openpyxl

    # "007" is not a canonical number, so it stays the string "007", not 7.
    assert openpyxl.load_workbook(wd / "data.xlsx").active["A2"].value == "007"


def test_edit_cell_text_flag_forces_string(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    # text=True keeps even a plain number as a string.
    tb.call("edit_cell", {"path": "data.xlsx", "cell": "B2", "value": "120", "text": True})

    import openpyxl

    cell = openpyxl.load_workbook(wd / "data.xlsx").active["B2"]
    assert cell.value == "120"
    assert isinstance(cell.value, str)


@pytest.mark.parametrize(
    ("tool_name", "arguments"),
    [
        ("read_xlsx", {}),
        ("edit_cell", {"cell": "A1", "value": "changed"}),
    ],
)
def test_invalid_xlsx_sheet_lists_available_sheets(tmp_path, tool_name, arguments):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    _add_xlsx_sheet(wd / "data.xlsx")

    out = tb.call(tool_name, {"path": "data.xlsx", "sheet": "Nope", **arguments})
    assert out == "error: no sheet 'Nope'; sheets: Sheet, Summary"


def test_read_and_edit_xlsx_explicit_sheet(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    _add_xlsx_sheet(wd / "data.xlsx")

    out = tb.call("read_xlsx", {"path": "data.xlsx", "sheet": "Summary"})
    assert "total" in out

    res = tb.call(
        "edit_cell",
        {"path": "data.xlsx", "sheet": "Summary", "cell": "B1", "value": "120"},
    )
    assert "100" in res and "120" in res

    import openpyxl

    assert openpyxl.load_workbook(wd / "data.xlsx")["Summary"]["B1"].value == 120


def test_xlsx_edit_is_undoable(tmp_path):
    tb, wd, rev = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    tb.call("edit_cell", {"path": "data.xlsx", "cell": "B2", "value": "999"})

    import openpyxl

    assert openpyxl.load_workbook(wd / "data.xlsx").active["B2"].value == 999
    rev.undo_last()  # restore the binary file exactly
    assert openpyxl.load_workbook(wd / "data.xlsx").active["B2"].value == 100


def test_read_and_edit_pptx(tmp_path):
    tb, wd, rev = _tb(tmp_path)
    _make_pptx(wd / "deck.pptx")

    out = tb.call("read_pptx", {"path": "deck.pptx"})
    assert "Old Title" in out

    res = tb.call(
        "edit_pptx_text",
        {"path": "deck.pptx", "find": "Old Title", "replace": "New Title"},
    )
    assert "1 run" in res

    from pptx import Presentation

    prs = Presentation(str(wd / "deck.pptx"))
    texts = [sh.text for sh in prs.slides[0].shapes if sh.has_text_frame]
    assert "New Title" in texts

    rev.undo_last()  # binary pptx restored
    prs2 = Presentation(str(wd / "deck.pptx"))
    texts2 = [sh.text for sh in prs2.slides[0].shapes if sh.has_text_frame]
    assert "Old Title" in texts2


def test_edit_pptx_missing_text_is_error(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_pptx(wd / "deck.pptx")
    res = tb.call("edit_pptx_text", {"path": "deck.pptx", "find": "nope", "replace": "x"})
    assert "not found" in res


def test_office_edit_outside_workspace_is_irreversible(tmp_path):
    """An office edit to a file outside the working dir isn't covered by the
    snapshot, so it's recorded as not-undoable (undo that doesn't lie)."""
    wd = tmp_path / "ws"
    wd.mkdir()
    rev = Reversibility(workdir=str(wd), rules=IgnoreRules())
    tb = Toolbox(str(wd), reversibility=rev, confirm=lambda p: True)

    outside = tmp_path / "outside.xlsx"  # sibling of ws/, not under it
    _make_xlsx(outside)

    res = tb.call("edit_cell", {"path": str(outside), "cell": "B2", "value": "999"})
    assert "999" in res
    last = rev.history()[-1]
    assert last.reversible is False
    assert "not undoable" in last.note


def test_office_edit_outside_declined_is_skipped(tmp_path):
    wd = tmp_path / "ws"
    wd.mkdir()
    rev = Reversibility(workdir=str(wd), rules=IgnoreRules())
    tb = Toolbox(str(wd), reversibility=rev, confirm=lambda p: False)  # decline

    outside = tmp_path / "outside.xlsx"
    _make_xlsx(outside)
    import openpyxl

    before = openpyxl.load_workbook(outside).active["B2"].value
    res = tb.call("edit_cell", {"path": str(outside), "cell": "B2", "value": "999"})
    assert res.startswith("skipped")
    after = openpyxl.load_workbook(outside).active["B2"].value
    assert after == before  # file on disk unchanged


def test_read_pptx_respects_max_slides(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    tb, wd, _ = _tb(tmp_path)
    prs = Presentation()
    for n in range(1, 6):  # 5 slides
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        shape.text_frame.text = f"Slide {n} body"
    prs.save(wd / "big.pptx")

    out = tb.call("read_pptx", {"path": "big.pptx", "max_slides": 3})
    assert "Slide 3 body" in out
    assert "Slide 4 body" not in out
    assert "... (2 more slides)" in out

    full = tb.call("read_pptx", {"path": "big.pptx"})  # default cap is 50
    assert "Slide 5 body" in full
    assert "more slides" not in full


def test_read_pptx_negative_cap_clamps_to_none(tmp_path):
    """A negative cap means 'show none', not garbage arithmetic (e.g. '6 more' for
    a 5-slide deck). It clamps to 0, so the count in the notice stays correct."""
    from pptx import Presentation
    from pptx.util import Inches

    tb, wd, _ = _tb(tmp_path)
    prs = Presentation()
    for n in range(1, 6):  # 5 slides
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        shape.text_frame.text = f"Slide {n} body"
    prs.save(wd / "big.pptx")

    out = tb.call("read_pptx", {"path": "big.pptx", "max_slides": -1})
    assert "Slide 1 body" not in out  # nothing shown
    assert "... (5 more slides)" in out  # correct count, not "6 more"


def test_read_docx_respects_max_paragraphs(tmp_path):
    pytest.importorskip("docx")
    from docx import Document

    tb, wd, _ = _tb(tmp_path)
    doc = Document()
    for n in range(1, 11):  # 10 paragraphs
        doc.add_paragraph(f"Paragraph {n}")
    doc.save(wd / "big.docx")

    out = tb.call("read_docx", {"path": "big.docx", "max_paragraphs": 4})
    assert "Paragraph 4" in out
    assert "Paragraph 5" not in out
    assert "... (6 more paragraphs)" in out

    full = tb.call("read_docx", {"path": "big.docx"})  # default cap is 200
    assert "Paragraph 10" in full
    assert "more paragraphs" not in full


# ---- append_rows / create_sheet ----


def test_append_rows_adds_after_the_last_row(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")  # header + one data row (alice)

    res = tb.call(
        "append_rows",
        {"path": "data.xlsx", "rows": [["bob", "80"], ["carol", "95"]]},
    )
    assert "appended 2 row(s) at 3-4" in res

    import openpyxl

    ws = openpyxl.load_workbook(wd / "data.xlsx").active
    assert ws.max_row == 4
    assert [ws["A3"].value, ws["B3"].value] == ["bob", 80]
    assert [ws["A4"].value, ws["B4"].value] == ["carol", 95]
    assert ws["A2"].value == "alice"  # existing rows untouched


def test_append_rows_coerces_like_edit_cell(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    tb.call(
        "append_rows",
        {"path": "data.xlsx", "rows": [["007", "1.5", "=B2*2", "hello"]]},
    )

    import openpyxl

    ws = openpyxl.load_workbook(wd / "data.xlsx").active
    assert ws["A3"].value == "007"  # leading zero preserved, as in edit_cell
    assert ws["B3"].value == 1.5
    assert ws["C3"].value == "=B2*2"  # formula passes through
    assert ws["D3"].value == "hello"


def test_append_rows_text_flag_keeps_everything_a_string(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    tb.call("append_rows", {"path": "data.xlsx", "rows": [["bob", "80"]], "text": True})

    import openpyxl

    cell = openpyxl.load_workbook(wd / "data.xlsx").active["B3"]
    assert cell.value == "80"
    assert isinstance(cell.value, str)


def test_append_rows_targets_a_named_sheet(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    _add_xlsx_sheet(wd / "data.xlsx")  # adds "Summary" with one row

    tb.call("append_rows", {"path": "data.xlsx", "rows": [["net", "42"]], "sheet": "Summary"})

    import openpyxl

    wb = openpyxl.load_workbook(wd / "data.xlsx")
    assert [wb["Summary"]["A2"].value, wb["Summary"]["B2"].value] == ["net", 42]
    assert wb.active.max_row == 2  # first sheet untouched


def test_append_rows_is_undoable(tmp_path):
    tb, wd, rev = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    before = (wd / "data.xlsx").read_bytes()

    tb.call("append_rows", {"path": "data.xlsx", "rows": [["bob", "80"]]})

    import openpyxl

    assert openpyxl.load_workbook(wd / "data.xlsx").active.max_row == 3
    rev.undo_last()  # restore the binary file exactly
    assert (wd / "data.xlsx").read_bytes() == before


def test_append_rows_rejects_a_flat_list(tmp_path):
    """["a", "b"] is the likely mistake; openpyxl would write it as two rows."""
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    before = (wd / "data.xlsx").read_bytes()

    res = tb.call("append_rows", {"path": "data.xlsx", "rows": ["bob", "80"]})
    assert "not a list" in res
    assert (wd / "data.xlsx").read_bytes() == before  # nothing written


def test_append_rows_empty_is_an_error(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    assert "non-empty list" in tb.call("append_rows", {"path": "data.xlsx", "rows": []})


def test_create_sheet_adds_a_named_sheet(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")

    res = tb.call("create_sheet", {"path": "data.xlsx", "name": "Summary"})
    assert "Summary" in res

    import openpyxl

    assert "Summary" in openpyxl.load_workbook(wd / "data.xlsx").sheetnames


def test_create_sheet_rejects_an_existing_name(tmp_path):
    """openpyxl would silently make 'Summary1'; the agent must be told instead."""
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    _add_xlsx_sheet(wd / "data.xlsx")

    res = tb.call("create_sheet", {"path": "data.xlsx", "name": "Summary"})
    assert "already exists" in res

    import openpyxl

    assert openpyxl.load_workbook(wd / "data.xlsx").sheetnames.count("Summary") == 1
    assert "Summary1" not in openpyxl.load_workbook(wd / "data.xlsx").sheetnames


def test_create_sheet_is_undoable(tmp_path):
    tb, wd, rev = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    before = (wd / "data.xlsx").read_bytes()

    tb.call("create_sheet", {"path": "data.xlsx", "name": "Summary"})

    import openpyxl

    assert "Summary" in openpyxl.load_workbook(wd / "data.xlsx").sheetnames
    rev.undo_last()
    assert (wd / "data.xlsx").read_bytes() == before


def test_append_rows_outside_workspace_is_irreversible(tmp_path):
    wd = tmp_path / "ws"
    wd.mkdir()
    rev = Reversibility(workdir=str(wd), rules=IgnoreRules())
    tb = Toolbox(str(wd), reversibility=rev, confirm=lambda p: True)

    outside = tmp_path / "outside.xlsx"
    _make_xlsx(outside)

    res = tb.call("append_rows", {"path": str(outside), "rows": [["bob", "80"]]})
    assert "appended 1 row(s)" in res
    last = rev.history()[-1]
    assert last.reversible is False
    assert "not undoable" in last.note


def test_append_rows_outside_declined_writes_nothing(tmp_path):
    wd = tmp_path / "ws"
    wd.mkdir()
    rev = Reversibility(workdir=str(wd), rules=IgnoreRules())
    tb = Toolbox(str(wd), reversibility=rev, confirm=lambda p: False)

    outside = tmp_path / "outside.xlsx"
    _make_xlsx(outside)
    before = outside.read_bytes()

    res = tb.call("append_rows", {"path": str(outside), "rows": [["bob", "80"]]})
    assert "skipped" in res
    assert outside.read_bytes() == before


def test_append_rows_to_a_fresh_sheet_starts_at_row_1(tmp_path):
    """openpyxl's ws.append() writes to row 2 on an empty sheet; we must not."""
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    tb.call("create_sheet", {"path": "data.xlsx", "name": "Blank"})

    res = tb.call(
        "append_rows",
        {"path": "data.xlsx", "rows": [["first", "1"]], "sheet": "Blank"},
    )
    assert "at 1-1" in res

    import openpyxl

    ws = openpyxl.load_workbook(wd / "data.xlsx")["Blank"]
    assert [ws["A1"].value, ws["B1"].value] == ["first", 1]
    assert ws.max_row == 1  # no stray blank row above it


def test_append_rows_twice_keeps_stacking(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    tb.call("append_rows", {"path": "data.xlsx", "rows": [["bob", "80"]]})
    res = tb.call("append_rows", {"path": "data.xlsx", "rows": [["carol", "95"]]})
    assert "at 4-4" in res

    import openpyxl

    ws = openpyxl.load_workbook(wd / "data.xlsx").active
    assert [ws["A3"].value, ws["A4"].value] == ["bob", "carol"]


def test_append_rows_handles_ragged_rows(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    tb.call(
        "append_rows",
        {"path": "data.xlsx", "rows": [["bob"], ["carol", "95", "extra"]]},
    )

    import openpyxl

    ws = openpyxl.load_workbook(wd / "data.xlsx").active
    assert [ws["A3"].value, ws["B3"].value] == ["bob", None]
    assert [ws["A4"].value, ws["B4"].value, ws["C4"].value] == ["carol", 95, "extra"]
