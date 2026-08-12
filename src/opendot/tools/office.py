"""Office-document tools: read/edit .xlsx and .pptx structurally.

These let the agent work with spreadsheets and presentations — file types that
`read_file`/`edit` can't touch (they're binary zips of XML). Each edit snapshots
the file first (via the same reversibility engine), so even a binary edit is
fully undoable: if the agent mangles a workbook, `undo` restores the exact bytes.
That "let the agent touch your spreadsheet, and cleanly revert if it botches it"
guarantee is unique to opendot.

Requires the ``office`` extra: ``pip install 'opendot[office]'``. If the libs
aren't installed, the tools aren't registered (see build_office_tools).
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

# A canonical, plain numeric literal: optional sign, ASCII digits, optional single
# decimal part and/or exponent. Deliberately strict — it rejects underscores
# ("1_000"), non-finite words ("inf"/"nan"), unicode digits, and thousands
# separators, all of which Python's int()/float() would otherwise accept. A
# leading-zero integer ("007") is handled separately so ZIP/ID/SKU values survive.
_INT_RE = re.compile(r"[+-]?[0-9]+")
_FLOAT_RE = re.compile(r"[+-]?(?:[0-9]+\.[0-9]*|\.[0-9]+|[0-9]+)(?:[eE][+-]?[0-9]+)?")


def _coerce_cell_value(value: str) -> Any:
    """Convert a cell string to int/float only when it's an unambiguous, canonical
    number; otherwise keep it as text. Formulas (leading '=') pass through as-is.

    Kept conservative on purpose: leading-zero integers, underscores, inf/nan, and
    non-ASCII digits all stay text, so IDs like "007" or codes aren't mangled.
    """
    if value.startswith("="):
        return value  # formula

    if _INT_RE.fullmatch(value):
        # Preserve leading zeros (ZIP/ID/SKU): "007" stays text, "0" is fine.
        digits = value.lstrip("+-")
        if len(digits) > 1 and digits[0] == "0":
            return value
        return int(value)

    if _FLOAT_RE.fullmatch(value):
        return float(value)  # regex already excludes inf/nan/underscores

    return value


def office_available() -> bool:
    try:
        import openpyxl  # noqa: F401
        import pptx  # noqa: F401

        return True
    except ImportError:
        return False


def build_office_tools(box) -> list:
    """Return the office Tools bound to a Toolbox (`box`). Empty if libs absent."""
    if not office_available():
        return []

    from opendot.tools.local import Tool, _truncate

    def _snapshot(p: Path) -> bool:
        """Snapshot ``p`` before an edit. Returns True to proceed, False if the
        user declined an out-of-workspace edit. A file outside the working dir
        isn't covered by the snapshot, so it can't be undone: confirm first and
        record it as irreversible, exactly like write_file/edit and the shell path.
        """
        outside = box._is_outside_workspace(p)
        if outside and not box._confirm(
            f"This edits a file outside the workspace and can't be undone:\n  {p}\nEdit it?"
        ):
            return False
        if box.rev is not None:
            box.rev.before_action(
                "write",
                str(p),
                reversible=not outside,
                note="outside the workspace — not undoable" if outside else "",
            )
        return True

    # ---- xlsx ----
    def read_xlsx(path: str, sheet: str | None = None, max_rows: int = 100) -> str:
        import openpyxl

        max_rows = max(0, max_rows)  # negative cap means "show none"
        p = box._resolve(path)
        if not p.exists():
            return f"error: file not found: {p}"
        wb = openpyxl.load_workbook(p, data_only=True)
        names = wb.sheetnames
        if sheet and sheet not in names:
            return f"error: no sheet {sheet!r}; sheets: {', '.join(names)}"
        ws = wb[sheet] if sheet else wb[names[0]]
        lines = [
            f"workbook {box._rel(p)}  sheets: {', '.join(names)}",
            f"sheet '{ws.title}'  ({ws.max_row} rows x {ws.max_column} cols)",
        ]
        for r, row in enumerate(ws.iter_rows(values_only=True), 1):
            if r > max_rows:
                lines.append(f"... ({ws.max_row - max_rows} more rows)")
                break
            cells = "\t".join("" if v is None else str(v) for v in row)
            lines.append(f"{r}: {cells}")
        return _truncate("\n".join(lines))

    def edit_cell(
        path: str, cell: str, value: str, sheet: str | None = None, text: bool = False
    ) -> str:
        import openpyxl

        p = box._resolve(path)
        if not p.exists():
            return f"error: file not found: {p}"
        wb = openpyxl.load_workbook(p)
        names = wb.sheetnames
        if sheet and sheet not in names:
            return f"error: no sheet {sheet!r}; sheets: {', '.join(names)}"
        ws = wb[sheet] if sheet else wb[names[0]]
        old = ws[cell].value
        if not _snapshot(p):
            return "skipped: user declined an edit outside the workspace"
        v = value if text else _coerce_cell_value(value)
        ws[cell] = v
        wb.save(p)
        return f"{box._rel(p)} [{ws.title}] {cell}: {old!r} → {v!r}"

    # ---- pptx ----
    def read_pptx(path: str, max_slides: int = 50) -> str:
        from pptx import Presentation

        max_slides = max(0, max_slides)  # negative cap means "show none"
        p = box._resolve(path)
        if not p.exists():
            return f"error: file not found: {p}"
        prs = Presentation(str(p))
        lines = [f"presentation {box._rel(p)}  ({len(prs.slides)} slides)"]
        for i, slide in enumerate(prs.slides, 1):
            if i > max_slides:
                lines.append(f"... ({len(prs.slides) - max_slides} more slides)")
                break
            texts = [
                sh.text.strip() for sh in slide.shapes if sh.has_text_frame and sh.text.strip()
            ]
            lines.append(f"— slide {i} —")
            lines.extend(f"   {t}" for t in texts)
        return _truncate("\n".join(lines))

    def edit_pptx_text(path: str, find: str, replace: str) -> str:
        from pptx import Presentation

        p = box._resolve(path)
        if not p.exists():
            return f"error: file not found: {p}"
        prs = Presentation(str(p))
        hits = 0
        for slide in prs.slides:
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        if find in run.text:
                            run.text = run.text.replace(find, replace)
                            hits += 1
        if hits == 0:
            return f"error: text {find!r} not found in {box._rel(p)} (no change)"
        if not _snapshot(p):  # in-memory edits above are dropped; disk file untouched
            return "skipped: user declined an edit outside the workspace"
        prs.save(str(p))
        return f"{box._rel(p)}: replaced {find!r} → {replace!r} in {hits} run(s)"

    # ---- docx ----
    def read_docx(path: str, max_paragraphs: int = 200) -> str:
        from docx import Document

        max_paragraphs = max(0, max_paragraphs)  # negative cap means "show none"
        p = box._resolve(path)
        if not p.exists():
            return f"error: file not found: {p}"
        doc = Document(str(p))
        lines = [f"document {box._rel(p)}  ({len(doc.paragraphs)} paragraphs)"]
        for i, para in enumerate(doc.paragraphs, 1):
            if i > max_paragraphs:
                lines.append(f"... ({len(doc.paragraphs) - max_paragraphs} more paragraphs)")
                break
            text = para.text.strip()
            if text:
                lines.append(text)
        return _truncate("\n".join(lines))

    tools = [
        Tool(
            "read_xlsx",
            "Read an .xlsx spreadsheet: lists sheets and prints a sheet's cells as rows.",
            {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "sheet": {"type": "string"},
                    "max_rows": {
                        "type": "integer",
                        "description": "Max rows to print per sheet (default 100).",
                    },
                },
                "required": ["path"],
            },
            read_xlsx,
        ),
        Tool(
            "edit_cell",
            "Set one cell in an .xlsx (e.g. cell 'B4'). Canonical numbers and "
            "formulas ('=...') are handled; pass text=true to force a string. Undoable.",
            {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "cell": {
                        "type": "string",
                        "description": "A1-style ref, e.g. 'B4'.",
                    },
                    "value": {"type": "string"},
                    "sheet": {"type": "string"},
                    "text": {
                        "type": "boolean",
                        "description": (
                            "Write the value as text without numeric coercion "
                            "(e.g. keep '007' or a phone number as a string). "
                            "Default false."
                        ),
                    },
                },
                "required": ["path", "cell", "value"],
            },
            edit_cell,
        ),
        Tool(
            "read_pptx",
            "Read a .pptx presentation: outputs each slide's text outline.",
            {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "max_slides": {
                        "type": "integer",
                        "description": "Max slides to output (default 50).",
                    },
                },
                "required": ["path"],
            },
            read_pptx,
        ),
        Tool(
            "edit_pptx_text",
            "Find-and-replace text across all slides of a .pptx. Undoable.",
            {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "find": {"type": "string"},
                    "replace": {"type": "string"},
                },
                "required": ["path", "find", "replace"],
            },
            edit_pptx_text,
        ),
    ]

    # .docx support is gated separately: python-docx is an independent optional
    # dep, so read_docx only registers when it's importable (openpyxl/pptx users
    # without python-docx still get the xlsx/pptx tools).
    try:
        import docx  # noqa: F401

        tools.append(
            Tool(
                "read_docx",
                "Read a .docx Word document: outputs its paragraph text.",
                {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "max_paragraphs": {
                            "type": "integer",
                            "description": "Max paragraphs to output (default 200).",
                        },
                    },
                    "required": ["path"],
                },
                read_docx,
            )
        )
    except ImportError:
        pass

    return tools
